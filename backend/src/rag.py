import hashlib
import chromadb

_embedder = None
_chroma = None


def _get_embedder():
    global _embedder
    if _embedder is None:
        from sentence_transformers import SentenceTransformer
        _embedder = SentenceTransformer('all-MiniLM-L6-v2')
    return _embedder


def _get_chroma():
    global _chroma
    if _chroma is None:
        _chroma = chromadb.PersistentClient(path="./chroma_db")
    return _chroma


def extract_text(filepath: str) -> str:
    lower = filepath.lower()
    if lower.endswith('.pdf'):
        import fitz
        doc = fitz.open(filepath)
        return "\n".join(page.get_text() for page in doc)
    elif lower.endswith('.docx'):
        from docx import Document
        doc = Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs)
    elif lower.endswith('.pptx'):
        from pptx import Presentation
        prs = Presentation(filepath)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n".join(texts)
    else:
        raise ValueError(f"Unsupported file type: {filepath}")


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    words = text.split()
    chunks = []
    step = chunk_size - overlap
    for i in range(0, len(words), step):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk:
            chunks.append(chunk)
    return chunks


def get_collection_name(user_id: int, note_id: int) -> str:
    key = f"user_{user_id}_note_{note_id}"
    return "rag_" + hashlib.md5(key.encode()).hexdigest()


def store_chunks(user_id: int, note_id: int, chunks: list[str]) -> None:
    name = get_collection_name(user_id, note_id)
    chroma = _get_chroma()
    try:
        chroma.delete_collection(name)
    except Exception:
        pass
    collection = chroma.create_collection(name)
    embeddings = _get_embedder().encode(chunks).tolist()
    collection.add(
        documents=chunks,
        embeddings=embeddings,
        ids=[f"chunk_{i}" for i in range(len(chunks))],
    )


def search_chunks(user_id: int, note_id: int, query: str, top_k: int = 5) -> list[str]:
    name = get_collection_name(user_id, note_id)
    collection = _get_chroma().get_collection(name)
    query_embedding = _get_embedder().encode([query]).tolist()
    results = collection.query(query_embeddings=query_embedding, n_results=top_k)
    return results['documents'][0] if results['documents'] else []


def delete_collection(user_id: int, note_id: int) -> None:
    name = get_collection_name(user_id, note_id)
    try:
        _get_chroma().delete_collection(name)
    except Exception:
        pass
