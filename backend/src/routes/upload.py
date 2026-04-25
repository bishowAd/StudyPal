from flask import Blueprint, request, jsonify
from flask_jwt_extended import get_jwt_identity
from src.middleware.auth import token_required
from src.models.subject import Subject
from src.models.topic import Topic
from src.ai import ask_ai_json
from extensions import db
import fitz
from docx import Document
from pptx import Presentation
import io
import json
import traceback

upload_bp = Blueprint('upload', __name__)

ALLOWED = {'.pdf', '.docx', '.pptx'}


def get_extension(filename: str) -> str:
    return '.' + filename.rsplit('.', 1)[-1].lower()


def extract_text(file, ext: str) -> str:
    if ext == '.pdf':
        doc  = fitz.open(stream=file.read(), filetype='pdf')
        text = ''.join(page.get_text() + '\n' for page in doc)
        doc.close()
        return text.strip()
    elif ext == '.docx':
        doc  = Document(io.BytesIO(file.read()))
        return '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
    elif ext == '.pptx':
        prs  = Presentation(io.BytesIO(file.read()))
        text = ''
        for i, slide in enumerate(prs.slides, 1):
            text += f'Slide {i}:\n'
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text += shape.text + '\n'
        return text.strip()
    return ''


# POST /api/upload/
@upload_bp.route('/', methods=['POST'])
@token_required
def upload():
    try:
        user_id = int(get_jwt_identity())  # cast to int — JWT stores as string
        file    = request.files.get('file')
        title   = request.form.get('title', '').strip()

        if not file:
            return jsonify({'message': 'No file uploaded'}), 400

        ext = get_extension(file.filename)
        if ext not in ALLOWED:
            return jsonify({'message': 'Only PDF, Word and PowerPoint supported'}), 400

        extracted = extract_text(file, ext)
        if not extracted or len(extracted) < 50:
            return jsonify({'message': 'Could not extract text from file'}), 400

        if not title:
            title = file.filename

        # Ask AI to extract all topics
        raw = ask_ai_json(
            prompt=(
                f'Extract all key topics and concepts from these study notes.\n\n'
                f'NOTES:\n{extracted[:8000]}\n\n'
                f'Return JSON:\n'
                f'{{"topics": ["topic 1", "topic 2", "topic 3"]}}\n\n'
                f'Rules:\n'
                f'1. Extract every distinct concept, theory, formula, term\n'
                f'2. Each topic should be 2-6 words\n'
                f'3. Return 10-30 topics depending on content\n'
                f'4. Be specific — not "Introduction" but "Cell Membrane Structure"'
            ),
            system='You are an expert at analyzing study material and extracting key topics. Return only valid JSON.'
        )

        parsed      = json.loads(raw)
        topic_names = parsed.get('topics', [])

        if not topic_names:
            return jsonify({'message': 'Could not extract topics from your notes'}), 500

        # Save subject
        subject = Subject(
            user_id        = user_id,
            title          = title,
            extracted_text = extracted,
        )
        db.session.add(subject)
        db.session.flush()

        # Save topics
        for name in topic_names:
            db.session.add(Topic(
                subject_id = subject.id,
                name       = name,
                status     = 'unknown',
            ))

        db.session.commit()

        return jsonify({
            'message': 'File uploaded and topics extracted',
            'subject': subject.to_dict(),
            'topics':  topic_names,
        }), 201

    except json.JSONDecodeError as e:
        return jsonify({
            'message': 'AI returned invalid JSON. Try again.',
            'error':   str(e),
        }), 500
    except Exception as e:
        db.session.rollback()
        return jsonify({
            'message': str(e),
            'error':   traceback.format_exc(),
        }), 500


# GET /api/upload/
@upload_bp.route('/', methods=['GET'])
@token_required
def get_subjects():
    try:
        user_id  = int(get_jwt_identity())
        subjects = Subject.query.filter_by(user_id=user_id)\
            .order_by(Subject.created_at.desc()).all()
        return jsonify({'subjects': [s.to_dict() for s in subjects]})
    except Exception as e:
        return jsonify({'message': str(e), 'error': traceback.format_exc()}), 500